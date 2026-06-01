# -*- coding: utf-8 -*-
"""
WirkVektor — Master-Folien-Bibliothek (16:9).
Erzeugt eine .pptx mit einer großen Anzahl gestalteter Layout-Folien für
alle gängigen Zwecke. Jede Folie folgt dem Designsystem aus DESIGN.md.
Bedienung: Folie duplizieren, Inhalt ersetzen.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

import brand as B

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.normpath(os.path.join(HERE, "..", "..", "..", "07 Anhänge", "Website Grafiken"))
OUT = os.path.normpath(os.path.join(HERE, "..", "Vorlagen", "Präsentation", "WirkVektor Master-Folien.pptx"))

# --- Geometrie ---------------------------------------------------------------
EMU_IN = 914400
SW = 13.333          # Folienbreite (Zoll)
SH = 7.5             # Folienhöhe
M  = 0.92            # Außenrand
CW = SW - 2 * M      # Inhaltsbreite

def C(hexs):
    return RGBColor.from_string(hexs)

# ----------------------------------------------------------------------------
# Bausteine
# ----------------------------------------------------------------------------
def add_slide(prs, bg=B.OFF_WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C(bg)
    return slide

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line); sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def text(slide, x, y, w, h, runs, size=16, color=B.SLATE_MID, bold=False,
         font=B.FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=1.0, caps=False, tracking=None, wrap=True):
    """runs: str (eine Ebene) oder Liste von Absätzen (je str)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    paras = runs if isinstance(runs, (list, tuple)) else [runs]
    for i, ptxt in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        try:
            p.line_spacing = spacing
        except Exception:
            pass
        r = p.add_run()
        r.text = (ptxt.upper() if caps else ptxt)
        f = r.font
        f.size = Pt(size); f.bold = bold; f.name = font
        f.color.rgb = C(color)
    return tb

def eyebrow(slide, x, y, label, color=B.VECTOR_TEAL, w=6.0):
    # kleiner Teal-Marker + Label-Caps
    rect(slide, x, y + 0.04, 0.34, 0.10, fill=color)
    text(slide, x + 0.5, y - 0.06, w, 0.3, label, size=11, color=color,
         bold=True, font=B.FONT_BODY, caps=True)

def logo(slide, x, y, scale=1.0, on_dark=False, wordmark=True):
    """WV-Glyph (Navy-Quadrat mit WV) + Wortmarke, gemäß favicon.svg."""
    s = 0.42 * scale
    glyph_fill = B.WHITE if on_dark else B.NAVY_DEEP
    glyph_txt  = B.NAVY_DEEP if on_dark else B.WHITE
    g = rect(slide, x, y, s, s, fill=glyph_fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
    tf = g.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "WV"
    r.font.size = Pt(15 * scale); r.font.bold = True; r.font.name = B.FONT_DISPLAY
    r.font.color.rgb = C(glyph_txt)
    if wordmark:
        text(slide, x + s + 0.16, y - 0.04, 3.0, s + 0.1, "WirkVektor",
             size=int(17 * scale), color=(B.WHITE if on_dark else B.NAVY_DEEP),
             bold=True, font=B.FONT_DISPLAY, anchor=MSO_ANCHOR.MIDDLE)

def footer(slide, page, on_dark=False, total="00"):
    sub = B.SLATE_LIGHT if not on_dark else B.SLATE_LIGHT
    text(slide, M, SH - 0.62, 6.0, 0.3, "WirkVektor — " + B.COMPANY["claim"],
         size=9, color=sub, font=B.FONT_BODY)
    text(slide, SW - M - 2.0, SH - 0.62, 2.0, 0.3, str(page),
         size=9, color=sub, font=B.FONT_BODY, align=PP_ALIGN.RIGHT)
    rect(slide, M, SH - 0.74, CW, 0.012, fill=(B.SAFETY_BORDER if not on_dark else "26324A"))

def card(slide, x, y, w, h, fill=B.WHITE, border=B.SAFETY_BORDER):
    return rect(slide, x, y, w, h, fill=fill, line=border, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)

def arrow(slide, x, y, w=0.5, h=0.28, color=B.VECTOR_TEAL):
    rect(slide, x, y, w, h, fill=color, shape=MSO_SHAPE.CHEVRON)

# ============================================================================
# FOLIEN
# ============================================================================
def s_cover(prs):
    s = add_slide(prs, B.NAVY_DEEP)
    logo(s, M, 0.85, scale=1.15, on_dark=True)
    eyebrow(s, M, 3.1, B.COMPANY["tagline"], color=B.IMPACT_CYAN, w=10)
    text(s, M, 3.45, CW, 2.4,
         ["Die Architektur", "wirksamer KI."],
         size=54, color=B.WHITE, bold=True, font=B.FONT_DISPLAY, spacing=1.02)
    rect(s, M, 6.05, 1.1, 0.07, fill=B.IMPACT_CYAN)
    text(s, M, 6.25, CW, 0.4, B.COMPANY["claim"], size=16,
         color="C7D0DE", font=B.FONT_BODY)
    # dezenter Vektor-Akzent rechts
    arrow(s, SW - 2.6, 3.7, w=1.6, h=0.5, color=B.VECTOR_TEAL)
    return s

def s_cover_light(prs):
    s = add_slide(prs, B.OFF_WHITE)
    logo(s, M, 0.85, scale=1.1)
    eyebrow(s, M, 3.1, "Präsentationstitel", w=10)
    text(s, M, 3.45, CW, 2.2, ["Aussagekräftiger", "Titel der Präsentation"],
         size=50, color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY, spacing=1.02)
    text(s, M, 6.0, CW, 0.5, "Untertitel · Anlass · Datum", size=16,
         color=B.SLATE_MID, font=B.FONT_BODY)
    rect(s, 0, SH - 0.16, SW, 0.16, fill=B.VECTOR_TEAL)
    return s

def s_title_minimal(prs):
    s = add_slide(prs, B.WHITE)
    eyebrow(s, M, 2.7, "Kapitelüberschrift", w=10)
    text(s, M, 3.05, CW, 1.4, "Klarer Titel in einer Zeile", size=40,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    text(s, M, 4.25, CW - 3, 1.2,
         "Ein knapper, sachlicher Untertitel, der den Inhalt der folgenden Folien einordnet.",
         size=18, color=B.SLATE_MID, font=B.FONT_BODY, spacing=1.3)
    footer(s, "03")
    return s

def s_agenda(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Agenda")
    text(s, M, 1.15, CW, 1.0, "Inhalt", size=34, color=B.NAVY_DEEP, bold=True,
         font=B.FONT_DISPLAY)
    items = ["Ausgangslage und Ziel", "KI-Readiness im Überblick",
             "Priorisierte Use Cases", "Governance und EU AI Act",
             "Umsetzungsfahrplan", "Nächste Schritte"]
    y = 2.5
    for i, it in enumerate(items, 1):
        text(s, M, y, 0.8, 0.5, f"{i:02d}", size=22, color=B.VECTOR_TEAL,
             bold=True, font=B.FONT_DISPLAY, anchor=MSO_ANCHOR.MIDDLE)
        text(s, M + 0.9, y, CW - 0.9, 0.5, it, size=18, color=B.NAVY_DEEP,
             font=B.FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, M, y + 0.6, CW, 0.012, fill=B.SAFETY_BORDER)
        y += 0.72
    footer(s, "04")
    return s

def s_divider_navy(prs):
    s = add_slide(prs, B.NAVY_DEEP)
    text(s, M, 2.4, 4.0, 2.0, "01", size=120, color="1E2A44", bold=True,
         font=B.FONT_DISPLAY)
    eyebrow(s, M, 4.4, "Abschnitt", color=B.IMPACT_CYAN)
    text(s, M, 4.75, CW, 1.2, "Titel des Abschnitts", size=40, color=B.WHITE,
         bold=True, font=B.FONT_DISPLAY)
    rect(s, M, 6.1, 1.1, 0.07, fill=B.IMPACT_CYAN)
    return s

def s_divider_teal(prs):
    s = add_slide(prs, B.VECTOR_TEAL)
    eyebrow_x = M
    rect(s, M, 3.05, 0.34, 0.10, fill=B.WHITE)
    text(s, M + 0.5, 2.95, 8, 0.3, "ABSCHNITT", size=11, color=B.WHITE, bold=True,
         font=B.FONT_BODY, caps=True)
    text(s, M, 3.35, CW, 1.6, "Wirkung beginnt mit Klarheit.", size=44,
         color=B.WHITE, bold=True, font=B.FONT_DISPLAY)
    logo(s, M, SH - 1.1, scale=0.9, on_dark=True)
    return s

def s_headline_lead(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Ausgangslage")
    text(s, M, 1.15, CW, 1.2, "Eine prägnante Kernaussage als Überschrift",
         size=30, color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY, spacing=1.1)
    text(s, M, 2.7, CW - 2.5, 3.0,
         ["Der einleitende Absatz fasst den Kontext sachlich zusammen. Er bleibt "
          "konkret und belegbar, ohne Übertreibung.",
          "",
          "Ein zweiter Absatz vertieft den Gedanken oder benennt die Konsequenz "
          "für das Unternehmen."],
         size=18, color=B.SLATE_MID, font=B.FONT_BODY, spacing=1.4)
    footer(s, "07")
    return s

def s_two_col(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Zwei Spalten")
    text(s, M, 1.15, CW, 0.9, "Gegenüberstellung in zwei Spalten", size=30,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    colw = (CW - 0.6) / 2
    for i, (h, body) in enumerate([
        ("Heute", "Beschreibung der aktuellen Situation in zwei bis drei Sätzen. "
                   "Sachlich und konkret."),
        ("Mit WirkVektor", "Beschreibung des angestrebten Zustands. Der Nutzen wird "
                   "messbar und nachvollziehbar.")]):
        x = M + i * (colw + 0.6)
        rect(s, x, 2.5, 0.34, 0.10, fill=B.VECTOR_TEAL)
        text(s, x, 2.7, colw, 0.5, h, size=20, color=B.NAVY_DEEP, bold=True,
             font=B.FONT_DISPLAY)
        text(s, x, 3.3, colw, 2.5, body, size=16, color=B.SLATE_MID,
             font=B.FONT_BODY, spacing=1.4)
    footer(s, "08")
    return s

def s_three_cards(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Drei Karten")
    text(s, M, 1.15, CW, 0.9, "Drei Bausteine im Überblick", size=30,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    cards = [("Strategie", "Use Cases identifizieren und nach Wirkung priorisieren."),
             ("Governance", "EU AI Act, Sicherheit und Datenschutz von Anfang an."),
             ("Wirkung", "Piloten produktiv führen und den Nutzen messen.")]
    gap = 0.5
    cw = (CW - 2 * gap) / 3
    for i, (h, body) in enumerate(cards):
        x = M + i * (cw + gap)
        card(s, x, 2.5, cw, 3.0)
        rect(s, x + 0.4, 2.9, 0.5, 0.5, fill=B.TEAL_TINT,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
        rect(s, x + 0.55, 3.05, 0.2, 0.2, fill=B.VECTOR_TEAL)
        text(s, x + 0.4, 3.7, cw - 0.8, 0.5, h, size=20, color=B.NAVY_DEEP,
             bold=True, font=B.FONT_DISPLAY)
        text(s, x + 0.4, 4.25, cw - 0.8, 1.5, body, size=15, color=B.SLATE_MID,
             font=B.FONT_BODY, spacing=1.35)
    footer(s, "09")
    return s

def s_four_cards(prs):
    s = add_slide(prs, B.NAVY_DEEP)
    rect(s, M, 0.89, 0.34, 0.10, fill=B.IMPACT_CYAN)
    text(s, M + 0.5, 0.79, 8, 0.3, "ZIELGRUPPEN", size=11, color=B.IMPACT_CYAN,
         bold=True, font=B.FONT_BODY, caps=True)
    text(s, M, 1.15, CW, 0.9, "Vier Rollen, ein Ziel", size=30, color=B.WHITE,
         bold=True, font=B.FONT_DISPLAY)
    roles = [("Geschäftsführung", "Wirtschaftlicher Nutzen und Steuerung."),
             ("IT-Leitung", "Architektur, Integration, Betrieb."),
             ("Sicherheit & Datenschutz", "Compliance und Kontrolle."),
             ("Fachbereiche", "Use Cases aus der Praxis.")]
    gap = 0.4
    cw = (CW - 3 * gap) / 4
    for i, (h, body) in enumerate(roles):
        x = M + i * (cw + gap)
        rect(s, x, 2.6, cw, 2.8, fill=B.NAVY_RAISED, line="26324A",
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        text(s, x + 0.3, 2.9, cw - 0.6, 0.3, f"{i+1:02d}", size=13,
             color=B.IMPACT_CYAN, bold=True, font=B.FONT_BODY)
        text(s, x + 0.3, 3.35, cw - 0.6, 1.0, h, size=17, color=B.WHITE,
             bold=True, font=B.FONT_DISPLAY, spacing=1.1)
        text(s, x + 0.3, 4.55, cw - 0.6, 0.8, body, size=13,
             color="A9B4C7", font=B.FONT_BODY, spacing=1.3)
    footer(s, "10", on_dark=True)
    return s

def s_bullets(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Aufzählung")
    text(s, M, 1.15, CW, 0.9, "Kernpunkte mit Akzent-Marker", size=30,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    pts = ["Sachliche, belegbare Aussage als erster Punkt.",
           "Jeder Punkt steht für genau einen Gedanken.",
           "Der Teal-Marker ersetzt klassische Bullet-Punkte.",
           "Konsistente Abstände sorgen für ruhiges Schriftbild.",
           "Maximal fünf bis sechs Punkte pro Folie."]
    y = 2.6
    for p in pts:
        rect(s, M, y + 0.08, 0.16, 0.16, fill=B.VECTOR_TEAL)
        text(s, M + 0.4, y - 0.05, CW - 0.4, 0.5, p, size=17, color=B.NAVY_DEEP,
             font=B.FONT_BODY, anchor=MSO_ANCHOR.TOP)
        y += 0.62
    footer(s, "11")
    return s

def s_quote(prs):
    s = add_slide(prs, B.NAVY_DEEP)
    text(s, M - 0.1, 1.6, 2, 1.5, "“", size=120, color=B.VECTOR_TEAL, bold=True,
         font=B.FONT_DISPLAY)
    text(s, M, 2.8, CW, 2.2,
         "KI entfaltet erst dann Wirkung, wenn sie kontrolliert, "
         "regelkonform und messbar in die Praxis kommt.",
         size=30, color=B.WHITE, bold=False, font=B.FONT_DISPLAY, spacing=1.25)
    rect(s, M, 5.5, 0.5, 0.05, fill=B.IMPACT_CYAN)
    text(s, M + 0.7, 5.3, CW, 0.5, ["Sebastian Schucht",],
         size=15, color=B.IMPACT_CYAN, bold=True, font=B.FONT_BODY)
    text(s, M + 0.7, 5.62, CW, 0.4, "Gründer & Geschäftsführer, WirkVektor",
         size=13, color="A9B4C7", font=B.FONT_BODY)
    return s

def s_definition(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Begriff")
    text(s, M, 1.3, CW, 1.0, "KI-Governance", size=40, color=B.NAVY_DEEP,
         bold=True, font=B.FONT_DISPLAY)
    card(s, M, 2.8, CW, 2.4, fill=B.WHITE)
    rect(s, M, 2.8, 0.09, 2.4, fill=B.VECTOR_TEAL)
    text(s, M + 0.6, 3.2, CW - 1.2, 1.8,
         "Der organisatorische Rahmen aus Regeln, Rollen und Kontrollen, der "
         "sicherstellt, dass KI im Unternehmen regelkonform, sicher und "
         "nachvollziehbar eingesetzt wird — Grundlage für EU-AI-Act-Konformität.",
         size=19, color=B.SLATE_MID, font=B.FONT_BODY, spacing=1.45,
         anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "13")
    return s

def s_dos_donts(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Empfehlung")
    text(s, M, 1.15, CW, 0.9, "Was zählt — und was schadet", size=30,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    colw = (CW - 0.6) / 2
    # Do
    card(s, M, 2.5, colw, 3.0, fill=B.TEAL_TINT, border=B.VECTOR_TEAL)
    text(s, M + 0.4, 2.8, colw - 0.8, 0.4, "EMPFOHLEN", size=12, color=B.VECTOR_TEAL,
         bold=True, font=B.FONT_BODY, caps=True)
    for i, t in enumerate(["Use Cases nach Wirkung priorisieren",
                            "Governance früh verankern",
                            "Mitarbeitende befähigen",
                            "Nutzen messbar machen"]):
        text(s, M + 0.4, 3.35 + i*0.5, colw - 0.8, 0.4, "+  " + t, size=15,
             color=B.NAVY_DEEP, font=B.FONT_BODY)
    # Don't
    x2 = M + colw + 0.6
    card(s, x2, 2.5, colw, 3.0, fill=B.WHITE, border=B.SAFETY_BORDER_STRONG)
    text(s, x2 + 0.4, 2.8, colw - 0.8, 0.4, "VERMEIDEN", size=12, color=B.SLATE_MID,
         bold=True, font=B.FONT_BODY, caps=True)
    for i, t in enumerate(["Tools ohne Strategie einführen",
                            "Governance nachträglich anbauen",
                            "Teams außen vor lassen",
                            "Wirkung nicht überprüfen"]):
        text(s, x2 + 0.4, 3.35 + i*0.5, colw - 0.8, 0.4, "–  " + t, size=15,
             color=B.SLATE_MID, font=B.FONT_BODY)
    footer(s, "14")
    return s

def s_kpi(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Wirkung in Zahlen")
    text(s, M, 1.15, CW, 0.9, "Kennzahlen mit Fortschrittsbalken", size=30,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    stats = [("68 %", "der Unternehmen ohne KI-Governance", 0.68),
             ("3×", "schnellere Use-Case-Auswahl", 0.85),
             ("–40 %", "Aufwand durch klare Priorisierung", 0.4)]
    gap = 0.5
    cw = (CW - 2*gap) / 3
    for i, (val, lab, frac) in enumerate(stats):
        x = M + i*(cw+gap)
        text(s, x, 2.7, cw, 1.2, val, size=52, color=B.VECTOR_TEAL, bold=True,
             font=B.FONT_DISPLAY)
        rect(s, x, 4.0, cw, 0.13, fill=B.SAFETY_BORDER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        rect(s, x, 4.0, cw*frac, 0.13, fill=B.VECTOR_TEAL,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        text(s, x, 4.3, cw, 0.9, lab, size=14, color=B.SLATE_MID,
             font=B.FONT_BODY, spacing=1.3)
    text(s, M, 5.7, CW, 0.4, "Beispielwerte — vor Nutzung mit belegten Quellen ersetzen.",
         size=11, color=B.SLATE_LIGHT, font=B.FONT_BODY)
    footer(s, "15")
    return s

def s_process(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Vorgehen")
    text(s, M, 1.15, CW, 0.9, "Vier Schritte zum Ergebnis", size=30,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    steps = [("01", "Analyse"), ("02", "Konzept"), ("03", "Pilot"), ("04", "Skalierung")]
    gap = 0.5
    cw = (CW - 3*gap) / 4
    y = 3.0
    rect(s, M + 0.2, y + 0.85, CW - 0.4, 0.02, fill=B.SAFETY_BORDER)
    for i, (n, t) in enumerate(steps):
        x = M + i*(cw+gap)
        rect(s, x + cw/2 - 0.45, y, 0.9, 0.9, fill=B.NAVY_DEEP,
             shape=MSO_SHAPE.OVAL)
        text(s, x + cw/2 - 0.45, y, 0.9, 0.9, n, size=22, color=B.WHITE, bold=True,
             font=B.FONT_DISPLAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x, y + 1.15, cw, 0.5, t, size=17, color=B.NAVY_DEEP, bold=True,
             font=B.FONT_DISPLAY, align=PP_ALIGN.CENTER)
        text(s, x, y + 1.65, cw, 1.0,
             "Kurze Beschreibung des Schritts in einem Satz.",
             size=13, color=B.SLATE_MID, font=B.FONT_BODY, align=PP_ALIGN.CENTER,
             spacing=1.3)
        if i < 3:
            arrow(s, x + cw + gap/2 - 0.18, y + 0.32, w=0.32, h=0.26)
    footer(s, "16")
    return s

def s_methodik(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Methodik")
    text(s, M, 1.15, CW, 0.9, "Die 6 Phasen der WirkVektor-Methodik", size=28,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    gap = 0.35
    cw = (CW - 2*gap) / 3
    ch = 1.65
    for i, (n, t, d) in enumerate(B.METHODIK):
        col = i % 3; row = i // 3
        x = M + col*(cw+gap); y = 2.45 + row*(ch+0.3)
        card(s, x, y, cw, ch)
        text(s, x + 0.35, y + 0.25, 1.2, 0.5, n, size=24, color=B.VECTOR_TEAL,
             bold=True, font=B.FONT_DISPLAY)
        text(s, x + 1.1, y + 0.28, cw - 1.3, 0.5, t, size=17, color=B.NAVY_DEEP,
             bold=True, font=B.FONT_DISPLAY)
        text(s, x + 0.35, y + 0.85, cw - 0.6, 0.7, d, size=12.5,
             color=B.SLATE_MID, font=B.FONT_BODY, spacing=1.25)
    footer(s, "17")
    return s

def s_timeline(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Fahrplan")
    text(s, M, 1.15, CW, 0.9, "Umsetzung über vier Quartale", size=30,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    y = 3.4
    rect(s, M, y, CW, 0.03, fill=B.SAFETY_BORDER_STRONG)
    qs = [("Q1", "Readiness & Strategie", True),
          ("Q2", "Use-Case-Sprint", True),
          ("Q3", "Governance & Pilot", False),
          ("Q4", "Skalierung & Messung", False)]
    cw = CW / 4
    for i, (q, t, done) in enumerate(qs):
        x = M + i*cw
        col = B.VECTOR_TEAL if done else B.SLATE_LIGHT
        rect(s, x + 0.0, y - 0.1, 0.22, 0.22, fill=col, shape=MSO_SHAPE.OVAL)
        text(s, x, y - 0.75, cw, 0.4, q, size=18, color=B.NAVY_DEEP, bold=True,
             font=B.FONT_DISPLAY)
        text(s, x, y + 0.3, cw - 0.4, 0.8, t, size=14, color=B.SLATE_MID,
             font=B.FONT_BODY, spacing=1.25)
    footer(s, "18")
    return s

def s_matrix(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Priorisierung")
    text(s, M, 1.15, CW, 0.9, "Nutzen-Aufwand-Matrix", size=30, color=B.NAVY_DEEP,
         bold=True, font=B.FONT_DISPLAY)
    ox, oy, sz = M + 1.0, 2.5, 3.4
    rect(s, ox, oy, sz, sz, fill=B.WHITE, line=B.SAFETY_BORDER_STRONG)
    rect(s, ox, oy, sz/2, sz/2, fill=B.TEAL_TINT)  # oben links: hoher Nutzen/geringer Aufwand
    rect(s, ox + sz/2, oy, 0.012, sz, fill=B.SAFETY_BORDER)
    rect(s, ox, oy + sz/2, sz, 0.012, fill=B.SAFETY_BORDER)
    text(s, ox + 0.2, oy + 0.2, sz/2 - 0.4, 0.6, "Quick Wins", size=14,
         color=B.VECTOR_TEAL, bold=True, font=B.FONT_BODY)
    # Achsen
    text(s, ox - 0.95, oy, 0.9, sz, "NUTZEN →", size=10, color=B.SLATE_MID,
         bold=True, font=B.FONT_BODY, caps=True, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, ox, oy + sz + 0.15, sz, 0.4, "AUFWAND →", size=10, color=B.SLATE_MID,
         bold=True, font=B.FONT_BODY, caps=True, align=PP_ALIGN.CENTER)
    for (dx, dy, lab) in [(0.9, 0.8, "A"), (2.4, 1.2, "B"), (1.4, 2.4, "C"),
                          (2.7, 2.7, "D")]:
        rect(s, ox+dx, oy+dy, 0.34, 0.34, fill=B.NAVY_DEEP, shape=MSO_SHAPE.OVAL)
        text(s, ox+dx, oy+dy, 0.34, 0.34, lab, size=12, color=B.WHITE, bold=True,
             font=B.FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Legende rechts
    lx = ox + sz + 0.8
    text(s, lx, oy + 0.1, CW - (lx - M), 0.4, "Use Cases", size=13,
         color=B.SLATE_LIGHT, bold=True, font=B.FONT_BODY, caps=True)
    for i, t in enumerate(["A — Dokumenten-Assistent", "B — Angebotsgenerierung",
                           "C — Wissensdatenbank (RAG)", "D — Prozessautomatisierung"]):
        text(s, lx, oy + 0.6 + i*0.5, CW - (lx - M), 0.4, t, size=14,
             color=B.NAVY_DEEP, font=B.FONT_BODY)
    footer(s, "19")
    return s

def s_comparison(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Optionen")
    text(s, M, 1.15, CW, 0.9, "Zwei Pakete im Vergleich", size=30,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    colw = (CW - 0.6) / 2
    data = [("KI-Readiness-Check", B.WHITE, B.NAVY_DEEP, B.SAFETY_BORDER,
             ["Standortbestimmung", "2–3 Wochen", "Management-Bericht", "Roadmap-Skizze"]),
            ("Produktiver KI-Pilot", B.NAVY_DEEP, B.WHITE, B.NAVY_DEEP,
             ["Use Case live", "6–10 Wochen", "Governance inklusive", "Wirkungsmessung"])]
    for i, (h, bg, fg, bd, feats) in enumerate(data):
        x = M + i*(colw+0.6)
        rect(s, x, 2.5, colw, 3.4, fill=bg, line=bd,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        acc = B.VECTOR_TEAL if i == 0 else B.IMPACT_CYAN
        rect(s, x + 0.4, 2.85, 0.34, 0.10, fill=acc)
        text(s, x + 0.4, 3.05, colw-0.8, 0.5, h, size=20, color=fg, bold=True,
             font=B.FONT_DISPLAY)
        for j, f in enumerate(feats):
            text(s, x + 0.4, 3.75 + j*0.48, colw-0.8, 0.4, "•  " + f, size=15,
                 color=(B.SLATE_MID if i == 0 else "C7D0DE"), font=B.FONT_BODY)
    footer(s, "20")
    return s

def s_table(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Übersicht")
    text(s, M, 1.15, CW, 0.9, "Tabelle im Markenstil", size=30, color=B.NAVY_DEEP,
         bold=True, font=B.FONT_DISPLAY)
    rows, cols = 5, 4
    gx = slide_table = s.shapes.add_table(rows, cols, Inches(M), Inches(2.5),
                                          Inches(CW), Inches(3.0)).table
    headers = ["Use Case", "Nutzen", "Aufwand", "Priorität"]
    body = [["Dokumenten-Assistent", "Hoch", "Gering", "1"],
            ["Wissensdatenbank (RAG)", "Hoch", "Mittel", "2"],
            ["Angebotsgenerierung", "Mittel", "Mittel", "3"],
            ["Prozessautomatisierung", "Hoch", "Hoch", "4"]]
    for c in range(cols):
        cell = gx.cell(0, c); cell.text = headers[c]
        cell.fill.solid(); cell.fill.fore_color.rgb = C(B.NAVY_DEEP)
        _style_cell(cell, B.WHITE, bold=True, size=13)
    for r in range(1, rows):
        for c in range(cols):
            cell = gx.cell(r, c); cell.text = body[r-1][c]
            cell.fill.solid()
            cell.fill.fore_color.rgb = C(B.WHITE if r % 2 else B.OFF_WHITE)
            _style_cell(cell, B.NAVY_DEEP if c == 0 else B.SLATE_MID,
                        bold=(c == 0), size=12)
    footer(s, "21")
    return s

def _style_cell(cell, color, bold=False, size=12):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.15); cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold; r.font.name = B.FONT_BODY
            r.font.color.rgb = C(color)

def s_bar_chart(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Auswertung")
    text(s, M, 1.15, CW, 0.9, "Balkendiagramm", size=30, color=B.NAVY_DEEP,
         bold=True, font=B.FONT_DISPLAY)
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3", "Q4"]
    cd.add_series("Wirkung", (18, 34, 52, 71))
    gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(M),
                                Inches(2.4), Inches(CW), Inches(3.6), cd)
    ch = gframe.chart
    ch.has_legend = False
    plot = ch.plots[0]
    plot.gap_width = 80
    for series in plot.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = C(B.VECTOR_TEAL)
    footer(s, "22")
    return s

def s_donut(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Verteilung")
    text(s, M, 1.15, CW, 0.9, "Ring-/Kreisdiagramm", size=30, color=B.NAVY_DEEP,
         bold=True, font=B.FONT_DISPLAY)
    cd = CategoryChartData()
    cd.categories = ["Strategie", "Governance", "Befähigung", "Betrieb"]
    cd.add_series("Anteil", (35, 25, 20, 20))
    gframe = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(M), Inches(2.3),
                                Inches(6.0), Inches(3.8), cd)
    ch = gframe.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.RIGHT
    ch.legend.include_in_layout = False
    pts = ch.plots[0].series[0].points
    palette = [B.NAVY_DEEP, B.VECTOR_TEAL, B.IMPACT_CYAN, B.SLATE_LIGHT]
    for i, pt in enumerate(pts):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = C(palette[i])
    footer(s, "23")
    return s

def s_image_full(prs):
    s = add_slide(prs, B.NAVY_DEEP)
    rect(s, 0, 0, SW, SH, fill=B.NAVY_RAISED)
    text(s, 0, SH/2 - 0.5, SW, 1.0, "[ Vollflächiges Bild ]", size=18,
         color=B.SLATE_LIGHT, font=B.FONT_BODY, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    # Untertitel-Leiste unten
    rect(s, 0, SH - 1.5, SW, 1.5, fill=B.NAVY_DEEP)
    eyebrow(s, M, SH - 1.25, "Bildunterschrift", color=B.IMPACT_CYAN)
    text(s, M, SH - 0.95, CW, 0.6, "Aussage zum Bild in einer prägnanten Zeile",
         size=22, color=B.WHITE, bold=True, font=B.FONT_DISPLAY)
    return s

def s_image_text(prs):
    s = add_slide(prs, B.OFF_WHITE)
    half = SW/2
    rect(s, 0, 0, half, SH, fill=B.SAFETY_BORDER)
    text(s, 0, SH/2 - 0.4, half, 0.8, "[ Bild ]", size=16, color=B.SLATE_MID,
         font=B.FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx = half + 0.7
    eyebrow(s, tx, 2.0, "Bild + Text")
    text(s, tx, 2.35, SW - tx - M, 1.2, "Bild links, Inhalt rechts", size=28,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY, spacing=1.1)
    text(s, tx, 3.7, SW - tx - M, 2.2,
         "Dieses Layout kombiniert eine Bildfläche mit erläuterndem Text. "
         "Geeignet für Referenzen, Produktbilder oder Porträts.",
         size=16, color=B.SLATE_MID, font=B.FONT_BODY, spacing=1.4)
    footer(s, "25")
    return s

def s_pakete(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Leistungspakete")
    text(s, M, 1.15, CW, 0.9, "Fünf standardisierte Pakete", size=30,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    gap = 0.3
    cw = (CW - 4*gap) / 5
    for i, name in enumerate(B.PAKETE):
        x = M + i*(cw+gap)
        card(s, x, 2.6, cw, 2.9)
        rect(s, x, 2.6, cw, 0.09, fill=B.VECTOR_TEAL)
        text(s, x + 0.25, 2.95, cw-0.5, 0.4, f"0{i+1}", size=18,
             color=B.VECTOR_TEAL, bold=True, font=B.FONT_DISPLAY)
        text(s, x + 0.25, 3.5, cw-0.5, 1.6, name, size=15, color=B.NAVY_DEEP,
             bold=True, font=B.FONT_DISPLAY, spacing=1.15)
    footer(s, "27")
    return s

def s_team(prs):
    s = add_slide(prs, B.OFF_WHITE)
    eyebrow(s, M, 0.85, "Team")
    text(s, M, 1.15, CW, 0.9, "Menschen hinter WirkVektor", size=30,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    portrait = os.path.join(ASSETS, "sebastian-schucht.png")
    pw = 2.2
    if os.path.exists(portrait):
        s.shapes.add_picture(portrait, Inches(M), Inches(2.6), height=Inches(2.75))
    else:
        rect(s, M, 2.6, pw, 2.75, fill=B.SAFETY_BORDER)
    tx = M + pw + 0.6
    text(s, tx, 2.7, CW - pw - 0.6, 0.5, "Sebastian Schucht", size=22,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    text(s, tx, 3.2, CW - pw - 0.6, 0.4, "Gründer & Geschäftsführer", size=14,
         color=B.VECTOR_TEAL, bold=True, font=B.FONT_BODY)
    text(s, tx, 3.7, CW - pw - 0.6, 1.6,
         "Begleitet mittelständische Unternehmen bei der strukturierten "
         "Einführung von KI — vom ersten Use Case bis zum produktiven Betrieb. "
         "Hintergrund in Informationssicherheit und Datenschutz.",
         size=15, color=B.SLATE_MID, font=B.FONT_BODY, spacing=1.4)
    footer(s, "28")
    return s

def s_cta(prs):
    s = add_slide(prs, B.NAVY_DEEP)
    eyebrow(s, M, 2.2, "Nächster Schritt", color=B.IMPACT_CYAN, w=10)
    text(s, M, 2.55, CW, 1.6, ["Lassen Sie uns Ihren", "KI-Readiness-Check starten."],
         size=40, color=B.WHITE, bold=True, font=B.FONT_DISPLAY, spacing=1.05)
    # CTA-Button
    rect(s, M, 4.7, 3.0, 0.7, fill=B.VECTOR_TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.08)
    text(s, M, 4.7, 3.0, 0.7, "Erstgespräch vereinbaren", size=15, color=B.WHITE,
         bold=True, font=B.FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, M + 3.4, 4.7, 6, 0.7, B.COMPANY["email"] + "   ·   " + B.COMPANY["web"],
         size=14, color="C7D0DE", font=B.FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)
    return s

def s_thanks(prs):
    s = add_slide(prs, B.NAVY_DEEP)
    logo(s, M, 0.9, scale=1.0, on_dark=True)
    text(s, M, 3.0, CW, 1.4, "Vielen Dank.", size=54, color=B.WHITE, bold=True,
         font=B.FONT_DISPLAY)
    rect(s, M, 4.35, 1.1, 0.07, fill=B.IMPACT_CYAN)
    text(s, M, 4.7, CW, 1.5,
         [B.COMPANY["owner"] + " · " + B.COMPANY["role"],
          B.COMPANY["email"],
          B.COMPANY["web"] + "  ·  " + B.COMPANY["linkedin"]],
         size=15, color="C7D0DE", font=B.FONT_BODY, spacing=1.5)
    return s

def s_appendix(prs):
    s = add_slide(prs, B.OFF_WHITE)
    rect(s, 0, 0, 0.5, SH, fill=B.NAVY_DEEP)
    eyebrow(s, M, 3.2, "Anhang")
    text(s, M, 3.55, CW, 1.0, "Backup & weiterführende Folien", size=34,
         color=B.NAVY_DEEP, bold=True, font=B.FONT_DISPLAY)
    return s

# ============================================================================
SLIDES = [
    s_cover, s_cover_light, s_title_minimal, s_agenda, s_divider_navy,
    s_divider_teal, s_headline_lead, s_two_col, s_three_cards, s_four_cards,
    s_bullets, s_quote, s_definition, s_dos_donts, s_kpi, s_process,
    s_methodik, s_timeline, s_matrix, s_comparison, s_table, s_bar_chart,
    s_donut, s_image_full, s_image_text, s_pakete, s_team, s_cta, s_thanks,
    s_appendix,
]

def build():
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN))
    prs.slide_height = Emu(int(SH * EMU_IN))
    for fn in SLIDES:
        fn(prs)
    cp = prs.core_properties
    cp.title = "WirkVektor — Master-Folien"
    cp.author = "WirkVektor"
    cp.category = "Vorlage"
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"Gespeichert: {OUT}  ({len(SLIDES)} Folien)")

if __name__ == "__main__":
    build()
